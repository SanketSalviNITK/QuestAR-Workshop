from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Define color scheme (Dark theme feel)
    # Background: #0f172a, Text: #f8fafc, Primary: #6366f1

    def add_slide(title, content_list):
        slide_layout = prs.slide_layouts[1] # Bullet slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Style Title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Style Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = content_list[0] if content_list else ""
        
        for i in range(1, len(content_list)):
            p = tf.add_paragraph()
            p.text = content_list[i]
            p.level = 0

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "QuestAR Workshop"
    subtitle.text = "From Portfolio to Augmented Reality\nA One-Day Beginner Workshop"

    # Slide 2: Goals
    add_slide("Workshop Overview", [
        "Build a professional personal portfolio",
        "Introduction to Modern Web (HTML5, CSS3, JS)",
        "Learn 3D Graphics with Three.js",
        "Build a Web-based AR Treasure Hunt",
        "Deploy and Host live on GitHub and Vercel"
    ])

    # Slide 3: Session 1
    add_slide("Session 1: Your Digital Identity", [
        "The Core Web: HTML5 (Structure) and CSS3 (Style)",
        "Building a premium Portfolio from scratch",
        "Version Control with GitHub",
        "Cloud Hosting with Vercel (Continuous Deployment)"
    ])

    # Slide 4: Session 2
    add_slide("Session 2: Entering the 3D World", [
        "Web UI Design: Premium Glassmorphism aesthetics",
        "Introduction to Three.js and 3D Scenes",
        "Working with Geometries, Materials, and Lighting",
        "Animating your first 3D Gem"
    ])

    # Slide 5: Session 3
    add_slide("Session 3: Bridging Reality (AR)", [
        "What is Web AR? (Image Tracking vs SLAM)",
        "Integrating Mind.AR with Three.js",
        "Mapping 3D treasures to real-world markers",
        "Building game logic: Collection and Scoring"
    ])

    # Slide 6: Session 4
    add_slide("Session 4: The Creator's Lab", [
        "Open Activity: Group brainstorming and building",
        "Remixing the code with new models and markers",
        "Presentation and Global Launch",
        "Sharing your live URLs!"
    ])

    # Slide 7: Q&A
    add_slide("Next Steps & Resources", [
        "Continue learning: Three.js Journey, Mind.AR Docs",
        "Join the Community (Discord, GitHub)",
        "Q&A Session",
        "Thank you for building the future!"
    ])

    prs.save('QuestAR_Workshop_Presentation.pptx')
    print("Presentation created successfully: QuestAR_Workshop_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
