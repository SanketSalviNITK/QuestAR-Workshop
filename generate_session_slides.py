from pptx import Presentation
from pptx.util import Inches, Pt

def add_bullet_slide(prs, title, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = content_list[0] if content_list else ""
    for i in range(1, len(content_list)):
        p = tf.add_paragraph()
        p.text = content_list[i]

def create_session1():
    prs = Presentation()
    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Session 1: Your Digital Identity"
    slide.placeholders[1].text = "Web Development Basics & Hosting"
    
    add_bullet_slide(prs, "What is the Web?", ["HTML: Structure (The Skeleton)", "CSS: Style (The Skin)", "JavaScript: Logic (The Brain)"])
    add_bullet_slide(prs, "HTML Fundamentals", ["Tags, Elements, and Attributes", "Semantic HTML (<header>, <nav>, <section>)", "Links and Images"])
    add_bullet_slide(prs, "CSS Basics", ["Selectors and Properties", "The Box Model", "Flexbox and Layouts", "Google Fonts integration"])
    add_bullet_slide(prs, "Version Control: GitHub", ["Why use Git?", "Initializing a repository", "Pushing your code to the cloud"])
    add_bullet_slide(prs, "Vercel: Instant Hosting", ["Connecting GitHub to Vercel", "Automatic deployments", "Your first live URL!"])
    add_bullet_slide(prs, "Assignment 1", ["Create a personal portfolio", "Style it with modern CSS", "Host it live on Vercel"])
    
    prs.save('Session1_Your_Digital_Identity.pptx')

def create_session2():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Session 2: Entering the 3D World"
    slide.placeholders[1].text = "UI Design & Three.js Basics"
    
    add_bullet_slide(prs, "Modern UI Design", ["Glassmorphism: Blur and Transparency", "Premium color palettes", "Responsive game dashboards"])
    add_bullet_slide(prs, "Intro to Three.js", ["3D in the browser", "The WebGL API", "Installing/Importing Three.js"])
    add_bullet_slide(prs, "The Big Three", ["1. Scene: The 3D world", "2. Camera: Your eyes", "3. Renderer: The drawing engine"])
    add_bullet_slide(prs, "Geometries & Materials", ["Creating shapes (Box, Sphere, Octahedron)", "Materials (Colors, Shine, Metalness)"])
    add_bullet_slide(prs, "Lighting the Scene", ["Ambient Light vs Point Light", "Casting shadows and highlights"])
    add_bullet_slide(prs, "Animation Loop", ["requestAnimationFrame", "Rotating and moving 3D objects"])
    
    prs.save('Session2_Entering_the_3D_World.pptx')

def create_session3():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Session 3: Bridging Reality"
    slide.placeholders[1].text = "Augmented Reality with Mind.AR"
    
    add_bullet_slide(prs, "What is Augmented Reality?", ["Overlaying digital content on the real world", "Web-based AR vs Native Apps"])
    add_bullet_slide(prs, "Mind.AR Overview", ["Image Tracking fundamentals", "How markers work", "The Mind.AR Compiler tool"])
    add_bullet_slide(prs, "Mind.AR + Three.js", ["Setting up the AR controller", "Adding anchors to markers", "Connecting 3D assets"])
    add_bullet_slide(prs, "Game Logic", ["Detecting when a target is found", "Collection mechanics (Click to collect)", "Updating the Score UI"])
    add_bullet_slide(prs, "Mobile Testing", ["Permissions (Camera access)", "HTTPS requirements", "QR Code generation for fast testing"])
    
    prs.save('Session3_Bridging_Reality.pptx')

def create_session4():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Session 4: The Creator's Lab"
    slide.placeholders[1].text = "Open Activity & Project Launch"
    
    add_bullet_slide(prs, "Activity Overview", ["Collaborative building session", "Remixing the code for new ideas", "Rapid prototyping"])
    add_bullet_slide(prs, "Brainstorming", ["AR Business Cards", "Interactive Art Gallery", "Educational Scavenger Hunts"])
    add_bullet_slide(prs, "Customization Guide", ["Swapping 3D models", "Creating custom markers", "Adding sound effects"])
    add_bullet_slide(prs, "Deployment & Launch", ["Final git push", "Testing on multiple devices", "Sharing your URL with the world"])
    add_bullet_slide(prs, "Wrap-up", ["Next steps in your dev journey", "Portfolio update", "Certificate distribution!"])
    
    prs.save('Session4_Creators_Lab.pptx')

if __name__ == "__main__":
    create_session1()
    create_session2()
    create_session3()
    create_session4()
    print("Four session-wise presentations created successfully!")
